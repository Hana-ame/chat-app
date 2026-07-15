from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

ACCENT = RGBColor(0x58, 0x65, 0xF2)
DARK = RGBColor(0x1E, 0x1E, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xFA)
SUCCESS = RGBColor(0x23, 0xA5, 0x59)
WARN = RGBColor(0xF0, 0xA0, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tb

def add_para(text_frame, text, font_size=16, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p

def make_title_bar(slide, text, subtitle=None):
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.6), ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.9), text, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.5), subtitle, font_size=16, color=RGBColor(0xCC, 0xCC, 0xFF))

# ── Slide 1: 封面 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK)
# accent bar
add_shape(sl, Inches(0), Inches(3.2), Inches(0.15), Inches(1.6), ACCENT)
add_textbox(sl, Inches(1.2), Inches(2.8), Inches(10), Inches(1.2), "Chat App 开发纪实", font_size=44, bold=True, color=WHITE)
add_textbox(sl, Inches(1.2), Inches(3.9), Inches(10), Inches(0.7), "7 天从零到生产级全栈聊天应用", font_size=24, color=RGBColor(0xBB, 0xBB, 0xDD))
info_tb = add_textbox(sl, Inches(1.2), Inches(4.8), Inches(10), Inches(1.5), "", font_size=16, color=GRAY)
add_para(info_tb.text_frame, "Go + chi + SQLite  |  React + Vite + Zustand", 16, color=GRAY)
add_para(info_tb.text_frame, "30 轮迭代 · 318 次提交 · 80+ Bug 修复", 16, color=GRAY)
add_para(info_tb.text_frame, "2026-07-10 → 2026-07-16", 16, color=GRAY)

# ── Slide 2: 时间线 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "30 轮迭代时间线", "4 个阶段横跨 7 天")

phases = [
    ("Phase 1", "Mock 驱动前端迭代", "Day 1 · 第 1-9 轮", ACCENT),
    ("Phase 2", "真实后端对接", "Day 2 · 第 10-16 轮", WARN),
    ("Phase 3", "架构重构铺路", "Day 3-4 · 第 11-20 轮", SUCCESS),
    ("Phase 4", "测试 + 上锁收尾", "Day 5-7 · 第 21-30 轮", RGBColor(0xE0, 0x60, 0x60)),
]
for i, (phase, desc, days, color) in enumerate(phases):
    y = Inches(2.2) + Inches(i * 1.2)
    add_shape(sl, Inches(0.8), y, Inches(0.12), Inches(0.8), color)
    add_textbox(sl, Inches(1.2), y, Inches(1.5), Inches(0.4), phase, font_size=18, bold=True, color=color)
    add_textbox(sl, Inches(2.8), y, Inches(5), Inches(0.4), desc, font_size=18, color=DARK)
    add_textbox(sl, Inches(8), y, Inches(3), Inches(0.4), days, font_size=16, color=GRAY, alignment=PP_ALIGN.RIGHT)

add_textbox(sl, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8), "策略: 先 Mock 后端 → 快速验证前端 → 再对接真实后端 → 最后重构上锁", font_size=14, color=GRAY)

# ── Slide 3: Mock 策略 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "策略: 先飞起来 — MITM Mock 模式", "前端开发不受后端阻塞")

# left column
tb = add_textbox(sl, Inches(0.6), Inches(2.0), Inches(5.5), Inches(5), "", font_size=16, color=DARK)
tf = tb.text_frame
tf.paragraphs[0].text = "问题"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ACCENT
add_para(tf, "前端需等后端就绪才能开发 — 并行阻塞", 16, color=DARK)
add_para(tf, "", 8)
add_para(tf, "方案: MITM 模式 Mock", 20, bold=True, color=ACCENT)
add_para(tf, "• 拦截 API 调用，模拟完整数据层", 16)
add_para(tf, "• 前 9 轮全部在前端完成", 16)
add_para(tf, "• 3 小时修复 11 个 Bug", 16)
add_para(tf, "• 统一数据层: ensureData()", 16)
add_para(tf, "", 8)
add_para(tf, "经典案例 — AI 消息消失", 20, bold=True, color=WARN)
add_para(tf, "消息写入 → 即时 onMessageCreate → 可见", 16)
add_para(tf, "AI 回复: store 流式逐字 + 数据层持久化", 16)
add_para(tf, "轮询从数据层读 → AI 消息不丢失", 16)

# right column - code
code_bg = add_shape(sl, Inches(6.8), Inches(2.0), Inches(5.8), Inches(3.8), RGBColor(0x28, 0x2C, 0x34))
code = add_textbox(sl, Inches(7.0), Inches(2.2), Inches(5.4), Inches(3.4), "", font_size=11, color=RGBColor(0xAB, 0xB2, 0xBF))
ctf = code.text_frame
ctf.word_wrap = True
lines = [
    "// MITM Mock 核心: 消息双轨写入",
    "d.messages.push(userMsg);",
    "store.onMessageCreate(userMsg);    // → 即时",
    "",
    "// AI 消息: 流式逐字 + 数据持久",
    "const aiMsg = {",
    "  content: '', streaming: true,",
    "  source: async (emit) => {         // 逐字 emit",
    "    for (let i = 0; i < text.length; i++) {",
    "      emit(text[i]);",
    "      d.messages.find(m=>m.id===aiId).content = acc;",
    "    }",
    "  }",
    "};",
]
for i, line in enumerate(lines):
    if i == 0:
        ctf.paragraphs[0].text = line
        ctf.paragraphs[0].font.size = Pt(11)
        ctf.paragraphs[0].font.color.rgb = RGBColor(0xAB, 0xB2, 0xBF)
    else:
        add_para(ctf, line, 11, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1), space_after=Pt(0))

# ── Slide 4: Mock 产出 + 代价 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "前 9 轮 Mock 交付了什么", "速度的代价")

# deliverables
tb = add_textbox(sl, Inches(0.6), Inches(2.0), Inches(5.5), Inches(4.5), "", font_size=16, color=DARK)
tf = tb.text_frame
tf.paragraphs[0].text = "✅ 交付"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = SUCCESS
deliverables = [
    "消息流式发送 + AI Bot 逐字回复",
    "附件上传 / Reaction / 搜索 / 排序",
    "7 个 UI 组件、自动滚动、输入框伸缩",
    "红点未读 / Pinned / 头像上传 / Settings",
]
for d in deliverables:
    add_para(tf, "• " + d, 16, color=DARK)

# costs
tb2 = add_textbox(sl, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5), "", font_size=16, color=DARK)
tf2 = tb2.text_frame
tf2.paragraphs[0].text = "⚠ 代价"
tf2.paragraphs[0].font.size = Pt(20)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = WARN
costs = [
    "Mock 多用户支持拖到第 7 轮才修",
    "排序 comparator: undefined vs false 的 !! 陷阱",
    "Hover 菜单改 4 次才定案",
    "Pinned API 改名全量重做（6 文档）",
]
for c in costs:
    add_para(tf2, "• " + c, 16, color=DARK)

add_textbox(sl, Inches(0.6), Inches(6.2), Inches(11), Inches(0.5), "核心教训: API 契约文档应在第一天写好，Mock 应与后端数据形状一致", font_size=14, bold=True, color=WARN)

# ── Slide 5: 撞墙 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "撞墙: 真实后端对接全线崩溃", "16 项 API 差异、字段名不一致、404")

issues = [
    ("member_count=0", "字段名不一致", "群组大小显示异常"),
    ("chat.members 为空", "数据结构不匹配", "成员列表白屏"),
    ("Leave 变成 Delete", "API 路径混淆", "用户解散了群组"),
    ("头像不显示", "URL 格式不同", "用户信息不完整"),
    ("搜索无结果", "参数名不对齐", "功能完全失效"),
]
for i, (issue, cause, effect) in enumerate(issues):
    y = Inches(2.2) + Inches(i * 0.9)
    add_shape(sl, Inches(0.6), y, Inches(0.08), Inches(0.6), ACCENT)
    add_textbox(sl, Inches(1.0), y, Inches(2.5), Inches(0.5), issue, font_size=16, bold=True, color=ACCENT)
    add_textbox(sl, Inches(3.5), y, Inches(4), Inches(0.5), cause, font_size=16, color=DARK)
    add_textbox(sl, Inches(7.5), y, Inches(4.5), Inches(0.5), effect, font_size=16, color=RGBColor(0x88, 0x88, 0x88))

add_textbox(sl, Inches(0.6), Inches(6.2), Inches(11), Inches(0.5), "根因: Mock 与 Go 后端的数据形状不一致 — 先快后稳的代价", font_size=14, bold=True, color=WARN)

# ── Slide 6: Members 摇摆 + Reaction 三幕 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "架构摇摆中学习", "Members 存储 3 次尝试 · Reaction 从建到拆")

# Members 方案
tb = add_textbox(sl, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.5), "", font_size=16, color=DARK)
tf = tb.text_frame
tf.paragraphs[0].text = "Members 存储方案"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ACCENT
add_para(tf, "❌ Store map membersByChatId", 16, color=RGBColor(0xE0, 0x60, 0x60))
add_para(tf, "   → 轮询覆盖，弃用", 14, color=GRAY)
add_para(tf, "✅ 组件 useEffect + 本地 state", 16, color=SUCCESS)
add_para(tf, "   → 轮询已含完整 members", 14, color=GRAY)

# Reaction
tb2 = add_textbox(sl, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5), "", font_size=16, color=DARK)
tf2 = tb2.text_frame
tf2.paragraphs[0].text = "Reaction me 字段三幕剧"
tf2.paragraphs[0].font.size = Pt(20)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = ACCENT
add_para(tf2, "Act 1: Handler enrichReactions", 16, color=DARK)
add_para(tf2, "  N+1 解析 + WS 广播无效", 14, color=GRAY)
add_para(tf2, "Act 2: 专用 GET /:id/reactions", 16, color=DARK)
add_para(tf2, "  客户端自行查询 me 状态", 14, color=GRAY)
add_para(tf2, "Act 3: Store onReaction 计算 me", 16, color=SUCCESS)
add_para(tf2, "  WS 广播不含 me，客户端按 viewer 计算", 14, color=GRAY)

add_textbox(sl, Inches(0.6), Inches(6.2), Inches(11), Inches(0.5), "教训: 架构决策前应先验证「轮询怎么合并」、「广播给谁用」", font_size=14, bold=True, color=WARN)

# ── Slide 7: 问题 → Service 层 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "问题: Handler 直连 DB 的重复劳动", "每加一个 API 就要复制一遍权限检查 + 广播")

# before code
code_bg = add_shape(sl, Inches(0.6), Inches(2.0), Inches(5.8), Inches(1.8), RGBColor(0x28, 0x2C, 0x34))
tb = add_textbox(sl, Inches(0.8), Inches(2.2), Inches(5.4), Inches(1.4), "", font_size=12, color=RGBColor(0xAB, 0xB2, 0xBF))
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "// 之前: Handler 直调 s.DB"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0xE0, 0x60, 0x60)
add_para(tf, "func (h *Handler) Foo(w, r) {", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(2))
add_para(tf, "  // 权限检查重复 · Hub 广播重复 · 验证重复", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1))
add_para(tf, "  h.DB.DoSomething(...)", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1))
add_para(tf, "}", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1))

# after
code_bg2 = add_shape(sl, Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.6), RGBColor(0x28, 0x2C, 0x34))
tb2 = add_textbox(sl, Inches(0.8), Inches(4.4), Inches(5.4), Inches(2.2), "", font_size=12, color=RGBColor(0xAB, 0xB2, 0xBF))
tf2 = tb2.text_frame
tf2.word_wrap = True
tf2.paragraphs[0].text = "// 之后: Handler → Service → DB (+Hub)"
tf2.paragraphs[0].font.size = Pt(12)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = RGBColor(0x23, 0xA5, 0x59)
add_para(tf2, "func (s *ChatService) DoSomething(ctx, chatID, userID) {", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(2))
add_para(tf2, "  s.MustBeMember(ctx, chatID, userID)    // AuthZ 统一", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1))
add_para(tf2, "  result := s.DB.DoSomething(...)          // DB 操作", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1))
add_para(tf2, "  s.Hub.Broadcast(chatID, event)          // 广播统一", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1))
add_para(tf2, "  return result", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1))
add_para(tf2, "}", 12, color=RGBColor(0xAB, 0xB2, 0xBF), space_before=Pt(1))

# right - architecture comparison
tb3 = add_textbox(sl, Inches(7.0), Inches(2.0), Inches(5.8), Inches(2.0), "", font_size=18, color=DARK)
tf3 = tb3.text_frame
tf3.paragraphs[0].text = "架构变化"
tf3.paragraphs[0].font.size = Pt(22)
tf3.paragraphs[0].font.bold = True
tf3.paragraphs[0].font.color.rgb = ACCENT
add_para(tf3, "", 8)
add_para(tf3, "  之前:  Handler → DB", 16, bold=True, color=RGBColor(0xE0, 0x60, 0x60))
add_para(tf3, "", 8)
add_para(tf3, "  之后:  Handler → Service → DB", 16, bold=True, color=SUCCESS)
add_para(tf3, "                          ↕", 16, color=DARK)
add_para(tf3, "                      Hub 广播", 16, bold=True, color=ACCENT)

add_textbox(sl, Inches(7.0), Inches(5.5), Inches(5.8), Inches(1.5), "6 个新服务文件: authz, chat, message, member, reaction, errors\nHandler 代码量减半（chat.go -179/+86, messages.go -199/+88）", font_size=14, color=DARK)

# ── Slide 8: Service 层设计 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "Service 层 4 个关键设计", "清晰的分层边界")

designs = [
    ("Sentinel 错误", "mapServiceError 消除 HTTP 硬编码", "service.ErrForbidden → 403\nservice.ErrNotFound → 404"),
    ("AuthZ 解耦", "MustBeMember 被 4 个 Service 共用", "Chat / Message / Member / Reaction\n统一在 authz.go 中"),
    ("广播集中", "Service 层内调 Hub，handler 不复用", "AddReaction/RemoveReaction 都走\ns.Hub.BroadcastBroadcast"),
    ("WithTx 占位", "预备跨表事务", "当前透传 DB，未来直接\n插入事务逻辑"),
]
for i, (title, desc, detail) in enumerate(designs):
    y = Inches(2.0) + Inches(i * 1.3)
    add_shape(sl, Inches(0.6), y, Inches(0.08), Inches(1.0), ACCENT)
    add_textbox(sl, Inches(1.0), y, Inches(2.5), Inches(0.4), title, font_size=18, bold=True, color=ACCENT)
    add_textbox(sl, Inches(1.0), y + Inches(0.35), Inches(2.5), Inches(0.5), desc, font_size=13, color=DARK)
    # code box for detail
    cb = add_shape(sl, Inches(3.8), y, Inches(8.8), Inches(1.0), RGBColor(0xF0, 0xF0, 0xF8))
    tb = add_textbox(sl, Inches(4.0), y + Inches(0.05), Inches(8.4), Inches(0.9), detail, font_size=12, color=DARK)

# ── Slide 9: Realtime 重构 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "实时连接重构: RealtimeCoordinator", "混乱的 connectWS/SSE/Polling → 状态机")

# before/after
tb = add_textbox(sl, Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.5), "", font_size=16, color=DARK)
tf = tb.text_frame
tf.paragraphs[0].text = "之前"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0xE0, 0x60, 0x60)
add_para(tf, "Store 中 4 个 connect/disconnect 方法", 16, color=DARK)
add_para(tf, "无状态管理 → 双通道、重连逻辑混乱", 16, color=DARK)

tb2 = add_textbox(sl, Inches(6.8), Inches(2.0), Inches(5.8), Inches(1.5), "", font_size=16, color=DARK)
tf2 = tb2.text_frame
tf2.paragraphs[0].text = "之后"
tf2.paragraphs[0].font.size = Pt(18)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = SUCCESS
add_para(tf2, "RealtimeCoordinator 状态机 singleton", 16, color=DARK)
add_para(tf2, "IDLE↔CONNECTING↔CONNECTED↔DISCONNECTING", 16, color=DARK)

# state machine diagram
states = ["IDLE", "CONNECTING", "CONNECTED", "DISCONNECTING"]
for i, s in enumerate(states):
    x = Inches(0.8 + i * 3.1)
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), Inches(2.6), Inches(0.7), )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT if s == "CONNECTED" else RGBColor(0xE0, 0xE0, 0xE8)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = s
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE if s == "CONNECTED" else DARK
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Calibri"
    tf.word_wrap = False

# arrows would be too complex, use text
add_textbox(sl, Inches(0.8), Inches(5.2), Inches(11), Inches(0.8), "状态守卫锁防双通道 | _closeGuard 阻止重连 | transport onClose → 3s 后重连", font_size=14, color=GRAY)

# Mock 解耦
tb3 = add_textbox(sl, Inches(0.6), Inches(5.8), Inches(11), Inches(1.2), "", font_size=16, color=DARK)
tf3 = tb3.text_frame
tf3.paragraphs[0].text = "同步: Mock 解耦 — MOCKABLE[] → Proxy(realApi, { get }) | 消除循环依赖 | 删除 dev/mock-ws.js"
tf3.paragraphs[0].font.size = Pt(15)
tf3.paragraphs[0].font.color.rgb = DARK

# ── Slide 10: 测试覆盖 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "从 0 到 92.6% 测试覆盖", "前 23 轮零测试 → 第 24-30 轮系统化覆盖")

# table-like layout
headers = ["轮次", "包", "覆盖率", "内容"]
data = [
    ("24", "config, service, ws", "起步", "核心路径"),
    ("25", "handler", "~60%", "核心 handler"),
    ("26", "service 错误路径", "↑", "取消上下文等边界"),
    ("30", "config / service / db", "100% / 92.6% / 81.6%", "错误路径全覆盖"),
]
for j, row in enumerate([headers] + data):
    y = Inches(2.2 + j * 0.6)
    for i, cell in enumerate(row):
        x = Inches(0.8 + i * 2.8)
        fs = 14 if j == 0 else 13
        b = j == 0
        add_textbox(sl, x, y, Inches(2.6), Inches(0.5), cell, font_size=fs, bold=b, color=DARK)

# DB 拆分
tb = add_textbox(sl, Inches(0.6), Inches(5.3), Inches(11), Inches(1.5), "", font_size=16, color=DARK)
tf = tb.text_frame
tf.paragraphs[0].text = "DB 包拆分"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = ACCENT
add_para(tf, "chats.go → 5 文件: chats / messages / members / reactions / refresh_tokens", 16, color=DARK)
add_para(tf, "Migration 陷阱: V001 ASCII V(86) < i(105) → 在 init.sql 前执行 → 改名 000__init.sql", 14, color=WARN)

# ── Slide 11: 数据总览 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "交付了什么", "功能矩阵 + 关键指标")

features = [
    ("实时消息 + AI Bot", "✅"),
    ("群组管理 + 角色权限", "✅"),
    ("Reaction + 附件上传", "✅"),
    ("搜索 + 公告 Pinned", "✅"),
    ("注册/登录 + JWT", "✅"),
    ("速率限制 + CSP", "✅"),
]
for i, (feat, status) in enumerate(features):
    x = Inches(0.8 + (i % 3) * 4.0)
    y = Inches(2.2 + (i // 3) * 1.0)
    add_shape(sl, x, y, Inches(3.5), Inches(0.7), RGBColor(0xF0, 0xF0, 0xFA))
    add_textbox(sl, x + Inches(0.3), y + Inches(0.1), Inches(2.8), Inches(0.5), f"{status} {feat}", font_size=15, color=DARK)

# metrics
metrics = [("7 天", "开发周期"), ("30 轮", "迭代轮次"), ("318", "Git 提交"), ("~80+", "修复 Bug"), ("92.6%", "测试覆盖")]
for i, (val, label) in enumerate(metrics):
    x = Inches(0.8 + i * 2.4)
    y = Inches(4.5)
    add_textbox(sl, x, y, Inches(2.0), Inches(0.6), val, font_size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(sl, x, y + Inches(0.6), Inches(2.0), Inches(0.4), label, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5), "前端 Bundle: 316 KB | 后端测试包: 14 | Go 后端 ~5000 行 + React 前端 ~3500 行", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 12: 经验教训 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, WHITE)
make_title_bar(sl, "经验与反思", "做得好的 · 可以更好的")

tb = add_textbox(sl, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.5), "", font_size=16, color=DARK)
tf = tb.text_frame
tf.paragraphs[0].text = "✅ 做得好的"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = SUCCESS
goods = [
    "Mock 先行 → 前端迭代不受后端阻塞",
    "Service 层提取 → handler 职责单一",
    "测试覆盖 0→92.6% → 重构的安全垫",
    "MITM Mock → 解决轮询覆盖根本问题",
]
for g in goods:
    add_para(tf, "• " + g, 16, color=DARK)

tb2 = add_textbox(sl, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5), "", font_size=16, color=DARK)
tf2 = tb2.text_frame
tf2.paragraphs[0].text = "🔧 可以更好的"
tf2.paragraphs[0].font.size = Pt(20)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = WARN
betters = [
    "API 契约应在第 1 天对齐，而不是第 2 天",
    "Store 方案决策前应验证轮询怎么合并",
    '做了又改的 30% 工作可借设计文档避免',
    "Mock 多用户支持是 Day 1 就该做的事",
]
for b in betters:
    add_para(tf2, "• " + b, 16, color=DARK)

# ── Slide 13: 总结 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK)

# big arc
add_shape(sl, Inches(0), Inches(3.0), Inches(0.12), Inches(1.5), ACCENT)
add_textbox(sl, Inches(0.6), Inches(2.0), Inches(12), Inches(1.0), "7 天完整生命周期", font_size=36, bold=True, color=WHITE)

arc = "原型 ── 撞墙 ── 重构 ── 上锁"
add_textbox(sl, Inches(0.6), Inches(3.0), Inches(12), Inches(0.8), arc, font_size=28, bold=True, color=RGBColor(0xBB, 0xBB, 0xDD))

stages = [
    ("原型", "快速验证", ACCENT),
    ("撞墙", "暴露问题", WARN),
    ("重构", "夯实架构", SUCCESS),
    ("上锁", "守住质量", RGBColor(0xE0, 0x60, 0x60)),
]
for i, (stage, desc, color) in enumerate(stages):
    y = Inches(4.0 + i * 0.7)
    add_shape(sl, Inches(0.6), y, Inches(0.08), Inches(0.5), color)
    add_textbox(sl, Inches(1.0), y, Inches(1.5), Inches(0.4), stage, font_size=18, bold=True, color=color)
    add_textbox(sl, Inches(2.8), y, Inches(3), Inches(0.4), desc, font_size=16, color=WHITE)

# quote
add_textbox(sl, Inches(0.6), Inches(5.8), Inches(12), Inches(1.2), '"先跑通，再重构，最后上锁 — 但别等到第 7 天才开始写测试。"', font_size=18, bold=False, color=GRAY)

# Save
output_path = "/mnt/d/WorkPlace/chat-app/docs/Chat-App-开发纪实.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"File size: {os.path.getsize(output_path) / 1024:.0f} KB")
